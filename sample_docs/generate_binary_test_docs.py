"""
Generate binary-format test documents for Sentinel E2E tests.
Requires: python-docx, openpyxl, python-pptx, fpdf

Run: python generate_binary_test_docs.py
"""
from pathlib import Path

OUT = Path(__file__).parent


# ── 1. DOCX — Employment contract, all 4 clauses ─────────────────────────────

def make_employment_docx():
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title = doc.add_heading("EMPLOYMENT AGREEMENT — CHIEF EXECUTIVE OFFICER", level=1)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(
        'This Employment Agreement ("Agreement") is entered into as of April 1, 2026, between '
        'NEXACORE SYSTEMS INC., a Delaware corporation ("Company"), and '
        'JONATHAN A. REED, an individual residing in Palo Alto, California ("Executive").'
    )

    # Section 1 — Compensation and Benefits
    doc.add_heading("SECTION 1 — COMPENSATION AND BENEFITS CLAUSE", level=2)
    doc.add_paragraph(
        "1.1  Base Salary: The Executive shall receive an annual base salary of USD 580,000, "
        "payable semi-monthly. The salary shall be reviewed annually by the Board of Directors."
    )
    doc.add_paragraph(
        "1.2  Annual Bonus: The Executive is eligible for an annual performance bonus of up to "
        "100% of base salary (target 60%), based on achievement of Board-approved KPIs."
    )
    doc.add_paragraph(
        "1.3  Equity Compensation: On the commencement date, the Executive shall receive options "
        "to purchase 1,200,000 shares of Common Stock at the fair market value on the grant date. "
        "Options vest over four (4) years: 25% after year one, then monthly thereafter."
    )
    doc.add_paragraph(
        "1.4  Benefits: The Company shall provide comprehensive health, dental, and vision "
        "insurance for Executive and dependants; life insurance of 3× base salary; 401(k) with "
        "4% employer match; USD 10,000 annual professional development allowance; and a car "
        "allowance of USD 1,500 per month."
    )
    doc.add_paragraph(
        "1.5  Paid Leave: The Executive shall be entitled to twenty-five (25) days of paid "
        "vacation per year, plus all Company-observed public holidays."
    )

    # Section 2 — IP Assignment
    doc.add_heading("SECTION 2 — INTELLECTUAL PROPERTY ASSIGNMENT CLAUSE", level=2)
    doc.add_paragraph(
        "2.1  Assignment: The Executive hereby irrevocably assigns to the Company all right, "
        "title, and interest in and to any inventions, improvements, developments, processes, "
        "software code, trade secrets, copyrightable works, or other intellectual property "
        "conceived, created, or developed by the Executive during the term of employment that "
        "relate to the Company's business or result from use of Company resources."
    )
    doc.add_paragraph(
        "2.2  Pre-existing IP: Executive discloses pre-existing intellectual property in "
        "Schedule A attached hereto. Such pre-existing IP is excluded from this assignment."
    )
    doc.add_paragraph(
        "2.3  Cooperation: Executive agrees to execute all documents and take all actions "
        "necessary to perfect, register, and enforce the Company's intellectual property rights."
    )
    doc.add_paragraph(
        "2.4  Works for Hire: All copyrightable works prepared by Executive within the scope "
        "of employment shall be deemed 'works made for hire' under 17 U.S.C. § 101."
    )

    # Section 3 — Termination and Severance
    doc.add_heading("SECTION 3 — TERMINATION AND SEVERANCE CLAUSE", level=2)
    doc.add_paragraph(
        "3.1  At-Will Employment: The Executive's employment is at-will and may be terminated "
        "by either party at any time, with or without cause, subject to the notice and severance "
        "provisions below."
    )
    doc.add_paragraph(
        "3.2  Termination by Company Without Cause: The Company may terminate Executive's "
        "employment without cause by providing ninety (90) days' prior written notice or, at "
        "the Company's election, paying ninety (90) days' base salary in lieu of notice."
    )
    doc.add_paragraph(
        "3.3  Severance Package: Upon termination without cause or resignation for Good Reason, "
        "the Executive shall receive: (a) 18 months' base salary continuation; (b) continued "
        "health benefits for 18 months or until re-employment with equivalent benefits, "
        "whichever is earlier; (c) acceleration of 50% of unvested equity."
    )
    doc.add_paragraph(
        "3.4  Termination for Cause: No severance is payable if Executive is terminated for "
        "cause, defined as: fraud, wilful misconduct, material breach of Agreement, conviction "
        "of a felony, or repeated failure to perform material duties after written warning."
    )
    doc.add_paragraph(
        "3.5  Resignation Without Good Reason: Executive must provide ninety (90) days' notice. "
        "No severance is payable."
    )

    # Section 4 — Non-compete and Confidentiality
    doc.add_heading("SECTION 4 — NON-COMPETE AND CONFIDENTIALITY CLAUSE", level=2)
    doc.add_paragraph(
        "4.1  Confidentiality: During and after employment, Executive shall maintain in strict "
        "confidence all Confidential Information, including business strategies, financial data, "
        "customer relationships, product roadmaps, and M&A plans."
    )
    doc.add_paragraph(
        "4.2  Non-Competition: For twenty-four (24) months following separation, Executive "
        "shall not accept a C-suite or VP-level position at any direct competitor of the Company "
        "in the enterprise software platform market within North America."
    )
    doc.add_paragraph(
        "4.3  Non-Solicitation of Employees: For twenty-four (24) months following separation, "
        "Executive shall not solicit, recruit, or hire any person who was an employee of the "
        "Company within twelve (12) months preceding separation."
    )
    doc.add_paragraph(
        "4.4  Non-Solicitation of Clients: For twenty-four (24) months following separation, "
        "Executive shall not solicit any client with whom the Company conducted business during "
        "the preceding two (2) years."
    )

    # Signatures
    doc.add_heading("SIGNATURES", level=2)
    doc.add_paragraph("NEXACORE SYSTEMS INC.\nBy: _________________________\nName: Dr. Anita M. Shah, Chairperson, Board of Directors\nDate: April 1, 2026")
    doc.add_paragraph("EXECUTIVE\nBy: _________________________\nName: Jonathan A. Reed\nDate: April 1, 2026")
    doc.add_paragraph("Agreement Reference: NCS-EA-CEO-2026-001")

    out_path = OUT / "employment_contract_ceo_word_format_all_clauses_APPROVED.docx"
    doc.save(str(out_path))
    print(f"  Created: {out_path.name}")


# ── 2. PDF — Credit agreement via fpdf ───────────────────────────────────────

def make_credit_pdf():
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "REVOLVING CREDIT FACILITY AGREEMENT", ln=True, align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.ln(4)

    pdf.multi_cell(0, 5,
        "This Revolving Credit Facility Agreement ('Agreement') is made as of February 20, 2026, "
        "between MIDWEST COMMERCIAL BANK N.A. ('Lender') and LAKEFRONT MANUFACTURING CORP., "
        "an Illinois corporation ('Borrower'). Facility Amount: USD 25,000,000.")
    pdf.ln(4)

    def section(title, body):
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(0, 8, title, ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5, body)
        pdf.ln(3)

    section(
        "SECTION 1 - GOVERNING LAW CLAUSE",
        "This Agreement shall be governed by and construed in accordance with the laws of the "
        "State of Illinois, without regard to its conflict-of-law principles. The Borrower "
        "irrevocably submits to the exclusive jurisdiction of the state and federal courts "
        "located in Cook County, Illinois for any dispute arising hereunder. The parties "
        "expressly waive any right to trial by jury in any proceeding arising from this "
        "Agreement. Any judgement of such courts shall be enforceable in any jurisdiction."
    )

    section(
        "SECTION 2 - EVENTS OF DEFAULT CLAUSE",
        "Each of the following constitutes an Event of Default: (a) Failure to pay any "
        "principal or interest within five (5) business days of the due date; (b) Material "
        "breach of any covenant, representation, or warranty not remedied within thirty (30) "
        "days of notice; (c) Commencement of insolvency, bankruptcy, or receivership proceedings; "
        "(d) Cross-default on indebtedness exceeding USD 2,000,000; (e) Material adverse change "
        "in Borrower's financial condition; (f) Change of control without prior written consent. "
        "Upon an Event of Default, Lender may declare all amounts immediately due and payable, "
        "cancel undrawn commitments, and exercise any security or other remedies available."
    )

    section(
        "SECTION 3 - INDEMNIFICATION CLAUSE",
        "Borrower shall indemnify, defend, and hold harmless Lender and its directors, officers, "
        "employees, advisors, and agents from any claims, losses, damages, costs, and expenses "
        "(including legal fees) arising from: (a) breach of this Agreement; (b) inaccuracy of "
        "any representation or warranty; (c) use of loan proceeds; (d) any environmental "
        "liability related to Borrower's properties; or (e) enforcement of this Agreement. "
        "The indemnification obligations shall survive repayment of the facility for three years."
    )

    section(
        "SECTION 4 - REPRESENTATIONS AND WARRANTIES",
        "Borrower represents and warrants that: (1) it is duly organized and in good standing "
        "under Illinois law; (2) it has full corporate authority to execute this Agreement; "
        "(3) financial statements provided are accurate and prepared in accordance with U.S. GAAP; "
        "(4) no material adverse change has occurred since the last audited balance sheet date; "
        "(5) Borrower is in compliance with all applicable laws and regulations; (6) no material "
        "litigation is pending or threatened; (7) all tax returns are current and taxes paid; "
        "(8) Borrower has good and marketable title to all collateral, free of liens. "
        "These representations and warranties are deemed repeated on each drawdown date."
    )

    pdf.set_font("Helvetica", "B", 11)
    pdf.cell(0, 8, "SIGNATURES", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 5,
        "MIDWEST COMMERCIAL BANK N.A.\nBy: _________________________\n"
        "Name: Thomas E. Kowalski, SVP Commercial Lending\nDate: February 20, 2026\n\n"
        "LAKEFRONT MANUFACTURING CORP.\nBy: _________________________\n"
        "Name: Patricia A. Novak, Chief Financial Officer\nDate: February 20, 2026\n\n"
        "Agreement Reference: MCB-RCF-LMC-2026-025")

    out_path = OUT / "credit_agreement_pdf_revolving_facility_APPROVED.pdf"
    pdf.output(str(out_path))
    print(f"  Created: {out_path.name}")


# ── 3. XLSX — Insurance policy ────────────────────────────────────────────────

def make_insurance_xlsx():
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Policy Summary"

    header_font = Font(bold=True, size=12, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    section_font = Font(bold=True, size=11, color="FFFFFF")
    section_fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin")
    )

    def header_cell(row, col, text):
        c = ws.cell(row=row, column=col, value=text)
        c.font = header_font
        c.fill = header_fill
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = thin_border

    def section_cell(row, col, text, colspan=1):
        c = ws.cell(row=row, column=col, value=text)
        c.font = section_font
        c.fill = section_fill
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        c.border = thin_border
        if colspan > 1:
            ws.merge_cells(start_row=row, start_column=col, end_row=row, end_column=col+colspan-1)

    def data_cell(row, col, text):
        c = ws.cell(row=row, column=col, value=text)
        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        c.border = thin_border

    ws.column_dimensions["A"].width = 35
    ws.column_dimensions["B"].width = 60
    ws.column_dimensions["C"].width = 25
    ws.row_dimensions[1].height = 30

    # Title row
    ws.merge_cells("A1:C1")
    title = ws.cell(row=1, column=1, value="COMMERCIAL PROPERTY AND CASUALTY INSURANCE POLICY")
    title.font = Font(bold=True, size=14, color="FFFFFF")
    title.fill = PatternFill(start_color="1A237E", end_color="1A237E", fill_type="solid")
    title.alignment = Alignment(horizontal="center", vertical="center")

    row = 2
    for label, value in [
        ("Policy Number", "CPC-2026-US-881234"),
        ("Insurer", "LIBERTY COMMERCIAL INSURANCE CO."),
        ("Insured", "LAKEFRONT MANUFACTURING CORP."),
        ("Policy Period", "March 1, 2026 to February 28, 2027"),
        ("Annual Premium", "USD 182,500"),
    ]:
        data_cell(row, 1, label)
        data_cell(row, 2, value)
        row += 1

    row += 1

    # Section 1 — Coverage Scope and Limits
    section_cell(row, 1, "SECTION 1 — COVERAGE SCOPE AND LIMITS", colspan=3)
    row += 1
    header_cell(row, 1, "Coverage Type")
    header_cell(row, 2, "Description")
    header_cell(row, 3, "Limit")
    row += 1
    coverages = [
        ("Commercial Property", "All-risk coverage for buildings, machinery, equipment, and inventory at 1450 Industrial Blvd, Chicago, IL. Replacement cost basis. Named perils exclusions apply.", "USD 15,000,000"),
        ("Business Interruption", "Loss of net income and continuing expenses following insured property damage. Includes extended period of indemnity of 180 days post-restoration.", "USD 5,000,000"),
        ("General Liability", "Bodily injury and property damage arising from business operations. Products and completed operations included. Defense costs outside the limit.", "USD 2,000,000 per occurrence / USD 4,000,000 aggregate"),
        ("Commercial Auto", "All owned, hired, and non-owned vehicles used for business. Medical payments, uninsured motorist, and liability coverage included.", "USD 1,000,000 combined single limit"),
        ("Workers Compensation", "Statutory benefits per Illinois Workers Compensation Act. Employer's Liability included. Coverage for all 847 employees.", "Statutory / USD 1,000,000 EL"),
        ("Inland Marine", "Tools, equipment, and goods in transit. Scheduled items per Exhibit A. Worldwide territorial scope during transit.", "USD 750,000"),
    ]
    for cov in coverages:
        for col, val in enumerate(cov, 1):
            data_cell(row, col, val)
        row += 1

    row += 1

    # Section 2 — Exclusions and Limitations
    section_cell(row, 1, "SECTION 2 — EXCLUSIONS AND LIMITATIONS CLAUSE", colspan=3)
    row += 1
    header_cell(row, 1, "Exclusion Category")
    header_cell(row, 2, "Description")
    header_cell(row, 3, "Standard/Endorseable")
    row += 1
    exclusions = [
        ("War and Terrorism", "Losses arising from war, invasion, civil war, terrorism, military action, or nuclear, biological, chemical, or radiological events.", "Standard"),
        ("Pollution", "Gradual or sudden pollution events unless sudden and accidental. Environmental cleanup costs excluded.", "Standard"),
        ("Cyber Events", "Data breach, ransomware, system failure losses unless Cyber endorsement (CYB-2026-001) is attached.", "Endorseable"),
        ("Wear and Tear", "Deterioration, rust, corrosion, mechanical breakdown without accompanying physical damage cause.", "Standard"),
        ("Intentional Acts", "Losses arising from intentional, fraudulent, or criminal acts by any insured.", "Standard"),
        ("Flood and Earthquake", "Flood and earthquake perils excluded. Separate NFIP and earthquake endorsements available.", "Endorseable"),
        ("Deductibles", "USD 10,000 per occurrence for property; USD 5,000 per occurrence for liability; USD 2,500 for auto.", "Standard"),
    ]
    for exc in exclusions:
        for col, val in enumerate(exc, 1):
            data_cell(row, col, val)
        row += 1

    row += 1

    # Section 3 — Claims Procedure
    section_cell(row, 1, "SECTION 3 — CLAIMS PROCEDURE", colspan=3)
    row += 1
    header_cell(row, 1, "Step")
    header_cell(row, 2, "Action Required")
    header_cell(row, 3, "Timeline")
    row += 1
    steps = [
        ("Step 1: Immediate Notice", "Notify Liberty Commercial Insurance at claims@libertycommercial.com or 1-800-555-CLAIM. Provide policy number, loss date, and brief description of incident.", "Within 24 hours of loss"),
        ("Step 2: Loss Mitigation", "Take all reasonable steps to prevent further loss. Document all emergency expenses. Do not authorize permanent repairs without insurer consent.", "Immediately upon loss"),
        ("Step 3: Formal Claim Filing", "Submit completed Proof of Loss form with supporting documentation: police report (if applicable), photos, inventory of damaged items, estimates.", "Within 30 days"),
        ("Step 4: Adjuster Assignment", "Insurer will assign a claims adjuster within 5 business days. Cooperate fully and provide access to damaged property.", "Within 5 business days"),
        ("Step 5: Settlement", "Insurer will provide settlement offer within 30 days of receiving complete claim documentation. Disputed amounts subject to appraisal.", "Within 30 days of docs"),
    ]
    for step in steps:
        for col, val in enumerate(step, 1):
            data_cell(row, col, val)
        row += 1

    row += 1

    # Section 4 — Premium Payment and Cancellation
    section_cell(row, 1, "SECTION 4 — PREMIUM PAYMENT AND CANCELLATION TERMS", colspan=3)
    row += 1
    header_cell(row, 1, "Item")
    header_cell(row, 2, "Details")
    header_cell(row, 3, "Amount / Date")
    row += 1
    premium_info = [
        ("Annual Premium (Total)", "Full premium for policy period March 1, 2026 to February 28, 2027", "USD 182,500"),
        ("Instalment 1 (25%)", "Due at policy inception", "USD 45,625 | Due: March 1, 2026"),
        ("Instalment 2 (25%)", "Second quarterly instalment", "USD 45,625 | Due: June 1, 2026"),
        ("Instalment 3 (25%)", "Third quarterly instalment", "USD 45,625 | Due: September 1, 2026"),
        ("Instalment 4 (25%)", "Fourth quarterly instalment", "USD 45,625 | Due: December 1, 2026"),
        ("Late Payment Fee", "Applies if instalment not received within 15 days of due date", "1.5% per month on overdue"),
        ("Cancellation by Insured", "30 days written notice required; pro-rata refund less 10% short-rate penalty", "Pro-rata refund"),
        ("Cancellation by Insurer", "10 days notice for non-payment; 30 days notice for other reasons; full pro-rata refund", "10-30 days notice"),
    ]
    for item in premium_info:
        for col, val in enumerate(item, 1):
            data_cell(row, col, val)
        row += 1

    out_path = OUT / "insurance_policy_xlsx_property_casualty_complete_APPROVED.xlsx"
    wb.save(str(out_path))
    print(f"  Created: {out_path.name}")


# ── 4. PPTX — Partnership agreement ──────────────────────────────────────────

def make_partnership_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    title_layout = prs.slide_layouts[0]  # title slide

    def add_text_slide(title_text, body_lines):
        slide = prs.slides.add_slide(blank_layout)

        # Title bar background
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)

        # Body text
        body_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.8))
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(body_lines):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = line
            p2.font.size = Pt(11)
            p2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            if line.startswith("   ") or line.startswith("•"):
                p2.level = 1

        return slide

    # Slide 1 — Cover (use add_text_slide helper for consistency)
    add_text_slide(
        "JOINT VENTURE PARTNERSHIP AGREEMENT",
        [
            "Parties: NovaTech Innovations Pte. Ltd. (Singapore) AND Atlantic Digital Ventures LLC (Delaware)",
            "",
            "Effective Date: March 1, 2026",
            "Agreement Reference: JV-NT-ADV-2026-001",
            "Governing Law: State of Delaware",
            "",
            "This agreement sets out the terms of a 50/50 joint venture for the development and",
            "commercialisation of AI-accelerated enterprise software in the Asia-Pacific and North American markets.",
        ]
    )

    # Slide 2 — Capital Contribution
    add_text_slide(
        "SECTION 1 — CAPITAL CONTRIBUTION CLAUSE",
        [
            "1.1  Initial Capital Contributions:",
            "   NovaTech Innovations Pte. Ltd. (Singapore): USD 8,000,000 cash + IP licensing rights valued at USD 4,000,000",
            "   Atlantic Digital Ventures LLC (Delaware): USD 10,000,000 cash + U.S. market distribution network valued at USD 2,000,000",
            "",
            "1.2  Ownership Split:  NovaTech — 50% | Atlantic Digital — 50%",
            "",
            "1.3  Additional Capital Calls: Either partner may propose a capital call by written notice.",
            "   Non-contributing partner's share diluted at the ratio of (existing capital / total post-call capital) × 100.",
            "",
            "1.4  Working Capital Reserve: Joint Venture shall maintain a minimum cash reserve equal to",
            "   3 months of projected operating expenses at all times.",
            "",
            "1.5  Capital Accounts: Each partner's capital account shall be maintained per U.S. GAAP.",
            "   Capital accounts adjusted annually for profit/loss allocations and any distributions taken.",
            "",
            "Reference: JV-NT-ADV-2026-001 | Capital Table recorded in Schedule A"
        ]
    )

    # Slide 3 — Profit and Loss Distribution
    add_text_slide(
        "SECTION 2 — PROFIT AND LOSS DISTRIBUTION",
        [
            "2.1  Profit Distribution:  Net profits shall be distributed in proportion to ownership interests:",
            "   NovaTech — 50% | Atlantic Digital — 50%",
            "",
            "2.2  Distribution Frequency: Distributions shall be made quarterly within 30 days of each",
            "   quarter-end, subject to the working capital reserve requirement of Section 1.4.",
            "",
            "2.3  Loss Allocation: Net losses are allocated in the same 50/50 proportion.",
            "   Partners may be required to make additional capital contributions to cover losses",
            "   that exceed existing capital account balances.",
            "",
            "2.4  Preferred Return: NovaTech receives a 6% cumulative preferred return on its IP",
            "   contribution (USD 4,000,000) before general profit sharing commences in Year 1.",
            "",
            "2.5  Reinvestment Policy: At least 20% of annual net profit shall be retained within",
            "   the Joint Venture for R&D investment, unless both partners agree otherwise.",
            "",
            "2.6  Audit Rights: Each partner may request an independent audit of JV financials",
            "   once per fiscal year at the requesting partner's expense."
        ]
    )

    # Slide 4 — Governance and Voting Rights
    add_text_slide(
        "SECTION 3 — GOVERNANCE AND VOTING RIGHTS",
        [
            "3.1  Board of Directors: The Joint Venture shall be governed by a Board of 6 Directors:",
            "   NovaTech appoints 3 directors | Atlantic Digital appoints 3 directors",
            "",
            "3.2  Chairperson: Rotating annually between NovaTech and Atlantic Digital nominees.",
            "   First Chairperson: NovaTech appointee for Year 1 (March 2026 – February 2027).",
            "",
            "3.3  Ordinary Resolutions: Approved by simple majority (4 of 6 directors).",
            "   Includes: operating budget approval, vendor contracts < USD 500,000.",
            "",
            "3.4  Supermajority Matters (5 of 6 directors required):",
            "   • Annual budget exceeding 15% above prior year",
            "   • Capital expenditures exceeding USD 2,000,000",
            "   • Entry into new geographic markets",
            "   • Incurring debt exceeding USD 5,000,000",
            "",
            "3.5  Unanimous Consent Required:",
            "   • Amendments to this Agreement",
            "   • Transfer of partnership interest to a third party",
            "   • Merger, acquisition, or dissolution of the Joint Venture",
        ]
    )

    # Slide 5 — Dissolution and Exit
    add_text_slide(
        "SECTION 4 — DISSOLUTION AND EXIT CLAUSE",
        [
            "4.1  Voluntary Dissolution: Either partner may propose dissolution by providing 180 days'",
            "   written notice. The Board must approve dissolution by unanimous consent.",
            "",
            "4.2  Winding Up Procedure: Upon dissolution, the Joint Venture shall:",
            "   (a) Continue ordinary operations to complete pending contracts",
            "   (b) Liquidate assets at fair market value",
            "   (c) Pay all creditors and third-party obligations",
            "   (d) Return capital contributions (adjusted for losses) to partners",
            "   (e) Distribute remaining assets 50/50 per ownership interests",
            "",
            "4.3  Buy-Out Option: Prior to dissolution, either partner may offer to purchase the",
            "   other's interest at fair market value (independent valuation required). The",
            "   non-initiating partner has 90 days to accept, decline, or counter-offer.",
            "",
            "4.4  Termination Events: Automatic dissolution triggered by: insolvency of either",
            "   partner, unanimous board resolution, or expiry of the 10-year initial term",
            "   (extendable by 5-year increments with unanimous consent).",
            "",
            "4.5  IP Reversion: Upon dissolution, all jointly developed IP vests equally.",
            "   NovaTech's licensed IP reverts to NovaTech."
        ]
    )

    # Slide 6 — Signatures
    add_text_slide(
        "EXECUTION — SIGNATURES",
        [
            "This Joint Venture Partnership Agreement has been duly executed by the authorised",
            "representatives of each party as of March 1, 2026.",
            "",
            "NOVATECH INNOVATIONS PTE. LTD.             ATLANTIC DIGITAL VENTURES LLC",
            "By: _________________________               By: _________________________",
            "Name: Dr. Wei Liang Chen                    Name: Michael T. O'Brien",
            "Title: Chief Executive Officer              Title: Managing Partner",
            "Date: March 1, 2026                        Date: March 1, 2026",
            "",
            "Agreement Reference: JV-NT-ADV-2026-001",
            "Governing Law: State of Delaware",
            "Arbitration: SIAC (Singapore International Arbitration Centre)",
        ]
    )

    out_path = OUT / "partnership_agreement_pptx_jv_tech_all_clauses_APPROVED.pptx"
    prs.save(str(out_path))
    print(f"  Created: {out_path.name}")


# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Generating binary test documents...")
    make_employment_docx()
    make_credit_pdf()
    make_insurance_xlsx()
    make_partnership_pptx()
    print("Done.")
