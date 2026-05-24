"""
Unit tests for data/file_extractor.py — unified multi-format text extraction.

TDD spec: extract_text(filename, content) dispatches by extension and returns
plain text. Each format is tested with a minimal in-memory fixture.
Run: pytest tests/unit/test_file_extractor.py -v
"""
import io
import pytest


# ---------------------------------------------------------------------------
# Minimal fixture factories (no disk files needed)
# ---------------------------------------------------------------------------

def _make_docx(text: str = "Force majeure clause applies here.") -> bytes:
    import docx
    doc = docx.Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(text: str = "Credit agreement governing law clause") -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = text
    ws["B1"] = "Events of default clause"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(text: str = "Dispute resolution clause arbitration") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = text
    slide.placeholders[1].text = "Limitation of liability details"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_html(text: str = "Force majeure clause is present.") -> bytes:
    return f"<html><head><title>Contract</title></head><body><h1>Agreement</h1><p>{text}</p></body></html>".encode()


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------

class TestDocxExtractor:
    def test_extracts_paragraph_text(self):
        from data.file_extractor import extract_text
        content = _make_docx("Limitation of liability shall not exceed fees paid.")
        result = extract_text("contract.docx", content)
        assert "Limitation of liability" in result

    def test_returns_string(self):
        from data.file_extractor import extract_text
        result = extract_text("doc.docx", _make_docx())
        assert isinstance(result, str)

    def test_strips_outer_whitespace(self):
        from data.file_extractor import extract_text
        result = extract_text("doc.docx", _make_docx("  some text  "))
        assert result == result.strip()

    def test_invalid_bytes_raise_value_error(self):
        from data.file_extractor import extract_text
        with pytest.raises(ValueError, match="valid DOCX"):
            extract_text("doc.docx", b"not a docx file at all")

    def test_multiple_paragraphs_joined(self):
        from data.file_extractor import extract_text
        import docx as _docx
        doc = _docx.Document()
        doc.add_paragraph("First paragraph.")
        doc.add_paragraph("Second paragraph.")
        buf = io.BytesIO()
        doc.save(buf)
        result = extract_text("multi.docx", buf.getvalue())
        assert "First paragraph" in result
        assert "Second paragraph" in result


# ---------------------------------------------------------------------------
# XLSX extraction
# ---------------------------------------------------------------------------

class TestXlsxExtractor:
    def test_extracts_cell_text(self):
        from data.file_extractor import extract_text
        result = extract_text("terms.xlsx", _make_xlsx("Governing law clause"))
        assert "Governing law clause" in result

    def test_returns_string(self):
        from data.file_extractor import extract_text
        result = extract_text("sheet.xlsx", _make_xlsx())
        assert isinstance(result, str)

    def test_multiple_cells_joined(self):
        from data.file_extractor import extract_text
        result = extract_text("terms.xlsx", _make_xlsx("Cell A1"))
        assert "Cell A1" in result
        assert "Events of default clause" in result

    def test_invalid_bytes_raise_value_error(self):
        from data.file_extractor import extract_text
        with pytest.raises(ValueError, match="valid XLSX"):
            extract_text("file.xlsx", b"this is not an xlsx")


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------

class TestPptxExtractor:
    def test_extracts_slide_text(self):
        from data.file_extractor import extract_text
        result = extract_text("deck.pptx", _make_pptx("Arbitration clause"))
        assert "Arbitration clause" in result

    def test_returns_string(self):
        from data.file_extractor import extract_text
        result = extract_text("deck.pptx", _make_pptx())
        assert isinstance(result, str)

    def test_multiple_shapes_joined(self):
        from data.file_extractor import extract_text
        result = extract_text("deck.pptx", _make_pptx("Title text"))
        assert "Title text" in result
        assert "Limitation of liability" in result

    def test_invalid_bytes_raise_value_error(self):
        from data.file_extractor import extract_text
        with pytest.raises(ValueError, match="valid PPTX"):
            extract_text("deck.pptx", b"not a pptx file")


# ---------------------------------------------------------------------------
# HTML extraction
# ---------------------------------------------------------------------------

class TestHtmlExtractor:
    def test_extracts_body_text(self):
        from data.file_extractor import extract_text
        result = extract_text("filing.html", _make_html("Indemnification clause applies."))
        assert "Indemnification clause applies" in result

    def test_strips_html_tags(self):
        from data.file_extractor import extract_text
        result = extract_text("page.html", b"<p>Clean <b>text</b> here.</p>")
        assert "<p>" not in result
        assert "<b>" not in result

    def test_excludes_script_content(self):
        from data.file_extractor import extract_text
        html = b"<html><body><p>Visible</p><script>alert('xss')</script></body></html>"
        result = extract_text("page.html", html)
        assert "alert" not in result
        assert "Visible" in result

    def test_excludes_style_content(self):
        from data.file_extractor import extract_text
        html = b"<html><head><style>.cls { color: red }</style></head><body><p>Content</p></body></html>"
        result = extract_text("page.html", html)
        assert "color" not in result
        assert "Content" in result

    def test_htm_extension_also_supported(self):
        from data.file_extractor import extract_text
        result = extract_text("doc.htm", _make_html("Force majeure"))
        assert "Force majeure" in result

    def test_returns_string(self):
        from data.file_extractor import extract_text
        result = extract_text("page.html", _make_html())
        assert isinstance(result, str)


# ---------------------------------------------------------------------------
# Image OCR extraction
# ---------------------------------------------------------------------------

class TestImageExtractor:
    def test_image_extraction_calls_tesseract(self, monkeypatch):
        from data import file_extractor
        monkeypatch.setattr(file_extractor, "_ocr_image", lambda content, _ext="": "Force majeure extracted via OCR")
        result = file_extractor.extract_text("scan.png", b"fake-image-bytes")
        assert "Force majeure" in result

    def test_jpg_extension_dispatches_to_ocr(self, monkeypatch):
        from data import file_extractor
        monkeypatch.setattr(file_extractor, "_ocr_image", lambda content, _ext="": "OCR result")
        result = file_extractor.extract_text("scan.jpg", b"fake")
        assert result == "OCR result"

    def test_tiff_extension_dispatches_to_ocr(self, monkeypatch):
        from data import file_extractor
        monkeypatch.setattr(file_extractor, "_ocr_image", lambda content, _ext="": "Tiff OCR text")
        result = file_extractor.extract_text("scan.tiff", b"fake")
        assert "Tiff OCR text" in result

    def test_ocr_unavailable_raises_value_error(self, monkeypatch):
        from data import file_extractor
        def _fail(content, _ext=""):
            raise ValueError("OCR not available: install Pillow and pytesseract.")
        monkeypatch.setattr(file_extractor, "_ocr_image", _fail)
        with pytest.raises(ValueError, match="OCR not available"):
            file_extractor.extract_text("scan.png", b"fake")


# ---------------------------------------------------------------------------
# Dispatcher — extension routing
# ---------------------------------------------------------------------------

class TestDispatcher:
    def test_pdf_dispatches_to_pdf_extractor(self, monkeypatch):
        from data import file_extractor
        monkeypatch.setattr(file_extractor, "_from_pdf", lambda c: "pdf text")
        result = file_extractor.extract_text("contract.pdf", b"fake")
        assert result == "pdf text"

    def test_txt_falls_back_to_utf8(self):
        from data.file_extractor import extract_text
        result = extract_text("notes.txt", b"Plain text document content.")
        assert "Plain text document content" in result

    def test_unknown_extension_tries_utf8(self):
        from data.file_extractor import extract_text
        result = extract_text("document.md", b"# Force Majeure\n\nClause text here.")
        assert "Force Majeure" in result

    def test_binary_non_text_raises_value_error(self):
        from data.file_extractor import extract_text
        with pytest.raises(ValueError):
            extract_text("blob.bin", bytes(range(256)))


# ---------------------------------------------------------------------------
# _safe_filename — path traversal protection (including Windows paths on Linux)
# ---------------------------------------------------------------------------

class TestSafeFilename:
    def test_strips_unix_directory(self):
        from data.file_extractor import _safe_filename
        assert _safe_filename("../../etc/passwd") == "passwd"

    def test_strips_windows_path_on_linux(self):
        from data.file_extractor import _safe_filename
        assert _safe_filename("C:\\Users\\evil\\file.pdf") == "file.pdf"

    def test_strips_mixed_slashes(self):
        from data.file_extractor import _safe_filename
        assert _safe_filename("path/to\\evil/../file.pdf") == "file.pdf"

    def test_plain_filename_unchanged(self):
        from data.file_extractor import _safe_filename
        assert _safe_filename("contract.pdf") == "contract.pdf"

    def test_rejects_dot(self):
        from data.file_extractor import _safe_filename
        with pytest.raises(ValueError, match="Invalid filename"):
            _safe_filename(".")

    def test_rejects_dotdot(self):
        from data.file_extractor import _safe_filename
        with pytest.raises(ValueError, match="Invalid filename"):
            _safe_filename("..")

    def test_rejects_empty_string(self):
        from data.file_extractor import _safe_filename
        with pytest.raises(ValueError, match="Invalid filename"):
            _safe_filename("")
