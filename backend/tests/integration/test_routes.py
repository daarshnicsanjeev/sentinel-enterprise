"""
Integration tests for FastAPI endpoints.

TDD spec: the HTTP boundary must validate inputs strictly and stream valid SSE.
The graph is mocked so these tests run without Ollama.
Run: pytest tests/integration/test_routes.py -v
"""
import json
import pytest
from fastapi.testclient import TestClient
from unittest.mock import AsyncMock, MagicMock, patch


@pytest.fixture(autouse=True)
def reset_rate_limiter():
    """Reset slowapi storage and doc_cache table before each test."""
    try:
        from api.routes import limiter
        limiter._storage.reset()
    except Exception:
        pass

    # Clear the doc_cache so dedup never returns cached results within tests
    try:
        import aiosqlite
        import asyncio
        import data.history_store as hs

        async def _clear_cache():
            async with aiosqlite.connect(hs._DB_PATH) as db:
                await db.execute("DELETE FROM doc_cache")
                await db.commit()

        asyncio.run(_clear_cache())
    except Exception:
        pass

    yield


@pytest.fixture
def client():
    from main import app
    return TestClient(app)


@pytest.fixture
def mock_graph_stream(monkeypatch):
    """Replace graph.astream with a deterministic async generator."""
    async def fake_astream(state, stream_mode=None):
        yield {
            "guardrail": {
                "sanitized": True,
                "logs": ["[Guardrail] Input sanitized: OK"],
                "final_decision": "PENDING",
            }
        }
        yield {
            "router": {
                "doc_type": "LEGAL_CONTRACT",
                "logs": ["[Router] Document classified as: LEGAL_CONTRACT"],
            }
        }
        yield {
            "compliance": {
                "final_decision": "REJECTED",
                "required_clauses": ["force majeure clause"],
                "compliance_output": "VERDICT: NON_COMPLIANT",
                "clause_results": [{"clause": "force majeure clause", "status": "MISSING"}],
                "logs": [
                    "[Compliance Tool] Queried regulatory DB for LEGAL_CONTRACT → Required: ['force majeure clause']",
                    "[Compliance] Verdict: REJECTED",
                ],
                "retry_count": 0,
            }
        }
        yield {
            "evaluator": {
                "evaluation_score": 0.85,
                "hallucination_risk": "low",
                "logs": ["[Evaluator] Faithfulness: 0.85 | Hallucination Risk: low | Accurate."],
            }
        }

    import agents.graph as graph_module
    graph_module.graph.astream = fake_astream


# ---------------------------------------------------------------------------
# Health endpoint
# ---------------------------------------------------------------------------

class TestHealthEndpoint:
    def test_returns_200(self, client):
        response = client.get("/api/health")
        assert response.status_code == 200

    def test_returns_ok_status(self, client):
        body = client.get("/api/health").json()
        assert body["status"] in ("ok", "degraded")

    def test_returns_service_name(self, client):
        body = client.get("/api/health").json()
        assert "service" in body
        assert "Sentinel" in body["service"]

    def test_returns_checks_dict(self, client):
        body = client.get("/api/health").json()
        assert "checks" in body
        assert isinstance(body["checks"], dict)

    def test_checks_include_database(self, client):
        body = client.get("/api/health").json()
        assert "database" in body["checks"]

    def test_checks_include_embeddings(self, client):
        body = client.get("/api/health").json()
        assert "embeddings" in body["checks"]

    def test_database_check_is_bool(self, client):
        body = client.get("/api/health").json()
        assert isinstance(body["checks"]["database"], bool)

    def test_ok_status_when_all_checks_pass(self, client):
        body = client.get("/api/health").json()
        checks = body["checks"]
        all_pass = all(v for v in checks.values())
        expected = "ok" if all_pass else "degraded"
        assert body["status"] == expected


# ---------------------------------------------------------------------------
# Analyze endpoint — input validation
# ---------------------------------------------------------------------------

class TestAnalyzeValidation:
    def test_no_file_returns_422(self, client):
        response = client.post("/api/analyze")
        assert response.status_code == 422

    def test_binary_file_returns_400(self, client, mock_graph_stream):
        binary_content = bytes(range(256))  # non-UTF-8 bytes
        files = {"file": ("test.bin", binary_content, "application/octet-stream")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 400

    def test_utf8_text_file_accepted(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200


# ---------------------------------------------------------------------------
# Analyze endpoint — SSE streaming
# ---------------------------------------------------------------------------

class TestAnalyzeStreaming:
    def test_response_content_type_is_event_stream(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert "text/event-stream" in response.headers["content-type"]

    def test_stream_contains_log_events(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        log_events = [p for p in payloads if p.get("type") == "log"]
        assert len(log_events) > 0

    def test_stream_ends_with_done_event(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done_events = [p for p in payloads if p.get("type") == "done"]
        assert len(done_events) == 1

    def test_done_event_has_final_decision(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done = next(p for p in payloads if p.get("type") == "done")
        assert "final_decision" in done

    def test_done_event_has_doc_type(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done = next(p for p in payloads if p.get("type") == "done")
        assert "doc_type" in done

    def test_log_events_have_node_and_message_fields(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document text for processing.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        log_events = [p for p in payloads if p.get("type") == "log"]
        for event in log_events:
            assert "node" in event
            assert "message" in event


# ---------------------------------------------------------------------------
# PDF upload path
# ---------------------------------------------------------------------------

def _make_minimal_pdf(text: str = "Force majeure clause applies to this agreement and governs both parties obligations.") -> bytes:
    """Minimal valid PDF for test uploads (PyMuPDF-generated, reliably extractable)."""
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 720), text)
    return doc.tobytes()


class TestPdfUpload:
    def test_valid_pdf_accepted(self, client, mock_graph_stream):
        files = {"file": ("contract.pdf", _make_minimal_pdf(), "application/pdf")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_pdf_response_is_event_stream(self, client, mock_graph_stream):
        files = {"file": ("contract.pdf", _make_minimal_pdf(), "application/pdf")}
        response = client.post("/api/analyze", files=files)
        assert "text/event-stream" in response.headers["content-type"]

    def test_invalid_pdf_bytes_returns_400(self, client, mock_graph_stream):
        files = {"file": ("bad.pdf", b"not a pdf", "application/pdf")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 400

    def test_pdf_stream_has_done_event(self, client, mock_graph_stream):
        files = {"file": ("contract.pdf", _make_minimal_pdf(), "application/pdf")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        assert any(p.get("type") == "done" for p in payloads)


# ---------------------------------------------------------------------------
# A1: File size guard
# ---------------------------------------------------------------------------

class TestFileSizeGuard:
    def test_rejects_file_over_5mb(self, client, mock_graph_stream):
        big_content = b"x" * (5 * 1024 * 1024 + 1)
        files = {"file": ("big.txt", big_content, "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 413

    def test_accepts_file_at_exact_5mb_limit(self, client, mock_graph_stream):
        ok_content = b"x" * (5 * 1024 * 1024)
        files = {"file": ("ok.txt", ok_content, "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_413_detail_mentions_size_limit(self, client, mock_graph_stream):
        big_content = b"x" * (5 * 1024 * 1024 + 1)
        files = {"file": ("big.txt", big_content, "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert "5 MB" in response.json()["detail"]


# ---------------------------------------------------------------------------
# A3: Request trace ID in SSE done event
# ---------------------------------------------------------------------------

class TestTraceId:
    def test_sse_done_event_includes_trace_id(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document for trace ID test.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done = next(p for p in payloads if p.get("type") == "done")
        assert "trace_id" in done

    def test_trace_id_is_uuid_format(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document for trace ID test.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done = next(p for p in payloads if p.get("type") == "done")
        # UUID4 format: 8-4-4-4-12 hex chars separated by hyphens
        import re
        assert re.match(
            r"^[0-9a-f]{8}-[0-9a-f]{4}-4[0-9a-f]{3}-[89ab][0-9a-f]{3}-[0-9a-f]{12}$",
            done["trace_id"],
        )


# ---------------------------------------------------------------------------
# B2: X-API-Key auth (optional — only enforced when SENTINEL_API_KEY is set)
# ---------------------------------------------------------------------------

class TestApiKeyAuth:
    def test_request_without_key_passes_when_no_key_configured(self, client, mock_graph_stream, monkeypatch):
        monkeypatch.delenv("SENTINEL_API_KEY", raising=False)
        import api.routes as routes_module
        monkeypatch.setattr(routes_module, "_API_KEY", "")
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_request_with_valid_key_passes(self, client, mock_graph_stream, monkeypatch):
        import api.routes as routes_module
        monkeypatch.setattr(routes_module, "_API_KEY", "test-secret")
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files, headers={"X-API-Key": "test-secret"})
        assert response.status_code == 200

    def test_request_with_wrong_key_returns_401(self, client, mock_graph_stream, monkeypatch):
        import api.routes as routes_module
        monkeypatch.setattr(routes_module, "_API_KEY", "test-secret")
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files, headers={"X-API-Key": "wrong-key"})
        assert response.status_code == 401

    def test_request_without_key_returns_401_when_key_configured(self, client, mock_graph_stream, monkeypatch):
        import api.routes as routes_module
        monkeypatch.setattr(routes_module, "_API_KEY", "test-secret")
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 401


# ---------------------------------------------------------------------------
# B4: Clause highlighting — done event includes clause_results
# ---------------------------------------------------------------------------

class TestClauseHighlighting:
    def test_sse_done_event_includes_clause_results(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done = next(p for p in payloads if p.get("type") == "done")
        assert "clause_results" in done

    def test_clause_results_is_list(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        lines = [l for l in response.text.split("\n") if l.startswith("data:")]
        payloads = [json.loads(l[5:].strip()) for l in lines]
        done = next(p for p in payloads if p.get("type") == "done")
        assert isinstance(done["clause_results"], list)


# ---------------------------------------------------------------------------
# C1: Persistent history
# ---------------------------------------------------------------------------

class TestHistory:
    def test_get_history_endpoint_exists(self, client):
        response = client.get("/api/history")
        assert response.status_code == 200

    def test_get_history_returns_list(self, client):
        response = client.get("/api/history")
        assert isinstance(response.json(), list)

    def test_analyze_persists_to_history(self, client, mock_graph_stream, monkeypatch):
        from unittest.mock import AsyncMock
        import api.routes as routes_module
        mock_insert = AsyncMock()
        monkeypatch.setattr(routes_module, "_save_to_history", mock_insert)
        files = {"file": ("contract.txt", b"Valid document text.", "text/plain")}
        client.post("/api/analyze", files=files)
        mock_insert.assert_called_once()


# ---------------------------------------------------------------------------
# C2: Human-in-the-loop override endpoint
# ---------------------------------------------------------------------------

class TestHITLOverride:
    def test_override_endpoint_exists_for_unknown_trace(self, client):
        response = client.post("/api/override/nonexistent-trace-id", json={"decision": "APPROVED"})
        # Should return 404 (trace not found) not 422 or 500
        assert response.status_code == 404

    def test_override_requires_decision_field(self, client):
        response = client.post("/api/override/some-id", json={})
        assert response.status_code == 422  # missing required field

    def test_override_rejects_invalid_decision(self, client):
        response = client.post("/api/override/some-id", json={"decision": "MAYBE"})
        assert response.status_code == 422


# ---------------------------------------------------------------------------
# P7.3 — CSV history export
# ---------------------------------------------------------------------------

class TestHistoryExport:
    def test_export_endpoint_returns_200(self, client):
        response = client.get("/api/history/export")
        assert response.status_code == 200

    def test_export_content_type_is_csv(self, client):
        response = client.get("/api/history/export")
        assert "text/csv" in response.headers["content-type"]

    def test_export_has_csv_header_row(self, client):
        response = client.get("/api/history/export")
        first_line = response.text.split("\n")[0]
        assert "trace_id" in first_line
        assert "decision" in first_line

    def test_export_content_disposition_header(self, client):
        response = client.get("/api/history/export")
        disposition = response.headers.get("content-disposition", "")
        assert "attachment" in disposition
        assert ".csv" in disposition

    def test_export_includes_history_records(self, client, monkeypatch):
        import data.history_store as hs
        async def fake_csv(limit=1000):
            return "trace_id,filename,doc_type,decision,faithfulness,risk,created_at\nabc,file.txt,LEGAL_CONTRACT,APPROVED,0.9,low,2026-01-01T00:00:00Z"
        monkeypatch.setattr(hs, "get_history_csv", fake_csv)
        response = client.get("/api/history/export")
        assert "abc" in response.text
        assert "APPROVED" in response.text


# ---------------------------------------------------------------------------
# P7.11 — Prometheus metrics endpoint
# ---------------------------------------------------------------------------

class TestMetricsEndpoint:
    def test_metrics_endpoint_returns_200(self, client):
        response = client.get("/api/metrics")
        assert response.status_code == 200

    def test_metrics_content_type_is_plaintext(self, client):
        response = client.get("/api/metrics")
        assert "text/plain" in response.headers["content-type"]

    def test_metrics_body_is_string(self, client):
        response = client.get("/api/metrics")
        assert isinstance(response.text, str)

    def test_metrics_incremented_after_analysis(self, client, mock_graph_stream, monkeypatch):
        import data.metrics as metrics_module
        monkeypatch.setattr(metrics_module, "_counters", {})
        files = {"file": ("doc.txt", b"Valid document text.", "text/plain")}
        client.post("/api/analyze", files=files)
        response = client.get("/api/metrics")
        assert "sentinel_analyses_total" in response.text


# ---------------------------------------------------------------------------
# P7.4 — Rate limiting
# ---------------------------------------------------------------------------

class TestRateLimit:
    def test_analyze_allows_request_within_limit(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_limiter_is_attached_to_app(self, client):
        from main import app
        # slowapi limiter must be registered on app.state
        assert hasattr(app.state, "limiter")


# ---------------------------------------------------------------------------
# P7.5 — Dead-letter queue
# ---------------------------------------------------------------------------

class TestDeadLetterQueue:
    def test_get_failures_endpoint_returns_200(self, client):
        response = client.get("/api/failures")
        assert response.status_code == 200

    def test_get_failures_returns_list(self, client):
        response = client.get("/api/failures")
        assert isinstance(response.json(), list)

    def test_failure_record_has_required_fields(self, client, monkeypatch):
        import data.history_store as hs
        async def fake_failures(limit=50):
            return [{"id": "f1", "trace_id": "t1", "filename": "test.txt", "error_msg": "Connection refused", "failed_at": "2026-05-17T00:00:00Z"}]
        monkeypatch.setattr(hs, "get_failures", fake_failures)
        response = client.get("/api/failures")
        failures = response.json()
        assert len(failures) == 1
        assert "trace_id" in failures[0]
        assert "error_msg" in failures[0]

    def test_pipeline_exception_stored_via_insert_failure(self, client, monkeypatch):
        """When _stream_graph raises, insert_failure should be called."""
        from unittest.mock import AsyncMock
        import api.routes as routes_module

        async def boom_stream(state, stream_mode=None):
            raise RuntimeError("Ollama died")
            yield  # make it a generator

        import agents.graph as graph_module
        graph_module.graph.astream = boom_stream

        mock_insert = AsyncMock()
        monkeypatch.setattr(routes_module, "_insert_failure", mock_insert)

        files = {"file": ("doc.txt", b"Valid document text.", "text/plain")}
        client.post("/api/analyze", files=files)
        mock_insert.assert_called_once()


# ---------------------------------------------------------------------------
# P7.9 — Webhook callback
# ---------------------------------------------------------------------------

class TestWebhook:
    def test_webhook_not_called_when_no_callback_url(self, client, mock_graph_stream, monkeypatch):
        from unittest.mock import AsyncMock
        import api.routes as routes_module
        mock_fire = AsyncMock()
        monkeypatch.setattr(routes_module, "_fire_webhook", mock_fire)
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        client.post("/api/analyze", files=files)
        mock_fire.assert_not_called()

    def test_invalid_url_scheme_returns_400(self, client, mock_graph_stream):
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        data = {"callback_url": "ftp://example.com/hook"}
        response = client.post("/api/analyze", files=files, data=data)
        assert response.status_code == 400

    def test_valid_callback_url_accepted(self, client, mock_graph_stream, monkeypatch):
        from unittest.mock import AsyncMock
        import api.routes as routes_module
        monkeypatch.setattr(routes_module, "_fire_webhook", AsyncMock())
        files = {"file": ("doc.txt", b"Valid document.", "text/plain")}
        data = {"callback_url": "https://example.com/hook"}
        response = client.post("/api/analyze", files=files, data=data)
        assert response.status_code == 200


# ---------------------------------------------------------------------------
# Multi-format file upload support
# ---------------------------------------------------------------------------

def _make_docx_bytes(text: str = "Force majeure clause in this agreement.") -> bytes:
    import io, docx
    doc = docx.Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes(text: str = "Governing law clause") -> bytes:
    import io, openpyxl
    wb = openpyxl.Workbook()
    wb.active["A1"] = text
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(text: str = "Dispute resolution slide") -> bytes:
    import io
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = text
    slide.placeholders[1].text = "Details here"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestMultiFormatUpload:
    def test_docx_upload_returns_200(self, client, mock_graph_stream):
        files = {"file": ("contract.docx", _make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_docx_response_is_event_stream(self, client, mock_graph_stream):
        files = {"file": ("contract.docx", _make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        response = client.post("/api/analyze", files=files)
        assert "text/event-stream" in response.headers["content-type"]

    def test_xlsx_upload_returns_200(self, client, mock_graph_stream):
        files = {"file": ("terms.xlsx", _make_xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_pptx_upload_returns_200(self, client, mock_graph_stream):
        files = {"file": ("deck.pptx", _make_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_html_upload_returns_200(self, client, mock_graph_stream):
        html = b"<html><body><p>Force majeure clause governs this agreement.</p></body></html>"
        files = {"file": ("filing.html", html, "text/html")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 200

    def test_invalid_docx_returns_400(self, client, mock_graph_stream):
        files = {"file": ("bad.docx", b"not a docx file", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 400

    def test_invalid_xlsx_returns_400(self, client, mock_graph_stream):
        files = {"file": ("bad.xlsx", b"not an xlsx file", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 400

    def test_invalid_pptx_returns_400(self, client, mock_graph_stream):
        files = {"file": ("bad.pptx", b"not a pptx file", "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        response = client.post("/api/analyze", files=files)
        assert response.status_code == 400
