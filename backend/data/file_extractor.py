"""
Unified multi-format text extractor for Project Sentinel.

Dispatches by file extension to format-specific extractors.
All extractors raise ValueError on invalid / unreadable input.
"""
import io
import os
import re
import zipfile
from pathlib import Path

# Point pytesseract at the Windows installer path if not already on PATH.
_TESSERACT_WIN = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if os.name == "nt" and os.path.exists(_TESSERACT_WIN):
    try:
        import pytesseract as _pt
        _pt.pytesseract.tesseract_cmd = _TESSERACT_WIN
    except ImportError:
        pass

# Maximum uncompressed size for ZIP-based formats (DOCX/XLSX/PPTX) — zip bomb protection
_MAX_ZIP_UNCOMPRESSED_BYTES = 50 * 1024 * 1024  # 50 MB

# Maximum extracted text length fed to agents — prevents FAISS OOM on text-heavy files
_MAX_EXTRACTED_CHARS = 200_000  # ~100 pages of text

# Magic byte signatures for format validation
_MAGIC_PDF  = b"%PDF"
_MAGIC_ZIP  = b"PK\x03\x04"   # DOCX / XLSX / PPTX are all ZIP containers
_MAGIC_PNG  = b"\x89PNG\r\n\x1a\n"
_MAGIC_JPEG = b"\xff\xd8\xff"
_MAGIC_TIFF_LE = b"II*\x00"   # little-endian TIFF
_MAGIC_TIFF_BE = b"MM\x00*"   # big-endian TIFF

# Maximum image pixels — explicit guard against PIL decompression bombs
# Default PIL limits are ~178M (warning) / ~357M (error); we set a tighter bound.
_MAX_IMAGE_PIXELS = 89_478_485  # 89M pixels (~9500×9500)


def _safe_filename(filename: str) -> str:
    """Return only the basename, rejecting path-traversal filenames.

    Normalises both forward and backward slashes so Windows-style paths like
    C:\\evil\\file.pdf are stripped correctly on Linux servers too.
    """
    normalized = filename.replace("\\", "/")
    safe = Path(normalized).name
    if not safe or safe in (".", ".."):
        raise ValueError("Invalid filename.")
    return safe


def _check_zip_bomb(content: bytes, fmt: str) -> None:
    """Reject ZIP-based documents whose uncompressed payload exceeds the limit."""
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            total = sum(info.file_size for info in zf.infolist())
            if total > _MAX_ZIP_UNCOMPRESSED_BYTES:
                mb = total // (1024 * 1024)
                raise ValueError(
                    f"{fmt} would decompress to {mb} MB, exceeding the 50 MB limit."
                )
    except zipfile.BadZipFile as exc:
        raise ValueError(f"not a valid {fmt}: {exc}") from exc


def _assert_magic(content: bytes, expected: bytes, fmt: str) -> None:
    """Verify file starts with expected magic bytes (guards against spoofed extensions)."""
    if not content[:len(expected)] == expected:
        raise ValueError(
            f"File does not appear to be a valid {fmt} (magic bytes mismatch)."
        )


def extract_text(filename: str, content: bytes) -> str:
    """Extract plain text from file bytes based on the file extension."""
    filename = _safe_filename(filename)
    ext = Path(filename.lower()).suffix

    if ext == ".pdf":
        text = _from_pdf(content)
    elif ext == ".docx":
        text = _from_docx(content)
    elif ext in (".xlsx", ".xls"):
        text = _from_excel(content)
    elif ext == ".pptx":
        text = _from_pptx(content)
    elif ext in (".html", ".htm"):
        text = _from_html(content)
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif"):
        text = _ocr_image(content, _ext=ext)
    else:
        text = _from_text(content)

    # Truncate to protect downstream FAISS indexing and LLM calls from huge inputs
    if len(text) > _MAX_EXTRACTED_CHARS:
        text = text[:_MAX_EXTRACTED_CHARS]

    return text


# ---------------------------------------------------------------------------
# Format-specific extractors (private)
# ---------------------------------------------------------------------------

def _from_pdf(content: bytes) -> str:
    _assert_magic(content, _MAGIC_PDF, "PDF")
    from data.pdf_extractor import extract_text_from_pdf
    return extract_text_from_pdf(content)


def _from_docx(content: bytes) -> str:
    _assert_magic(content, _MAGIC_ZIP, "DOCX")
    _check_zip_bomb(content, "DOCX")
    try:
        import docx
        doc = docx.Document(io.BytesIO(content))
    except zipfile.BadZipFile as exc:
        raise ValueError(f"not a valid DOCX: {exc}") from exc
    except Exception as exc:
        raise ValueError(f"not a valid DOCX: {exc}") from exc

    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("DOCX contains no extractable text")
    return text.strip()


def _from_excel(content: bytes) -> str:
    _assert_magic(content, _MAGIC_ZIP, "XLSX")
    _check_zip_bomb(content, "XLSX")
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    except zipfile.BadZipFile as exc:
        raise ValueError(f"not a valid XLSX: {exc}") from exc
    except Exception as exc:
        raise ValueError(f"not a valid XLSX: {exc}") from exc

    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(v) for v in row if v is not None and str(v).strip())
            if row_text:
                parts.append(row_text)
    wb.close()

    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("Excel file contains no extractable text")
    return text.strip()


def _from_pptx(content: bytes) -> str:
    _assert_magic(content, _MAGIC_ZIP, "PPTX")
    _check_zip_bomb(content, "PPTX")
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
    except zipfile.BadZipFile as exc:
        raise ValueError(f"not a valid PPTX: {exc}") from exc
    except Exception as exc:
        raise ValueError(f"not a valid PPTX: {exc}") from exc

    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

    text = "\n".join(parts)
    if not text.strip():
        raise ValueError("PPTX contains no extractable text")
    return text.strip()


def _from_html(content: bytes) -> str:
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self._parts: list[str] = []
            self._skip_tags = {"script", "style", "head"}
            self._skip_depth = 0

        def handle_starttag(self, tag, attrs):
            if tag in self._skip_tags:
                self._skip_depth += 1

        def handle_endtag(self, tag):
            if tag in self._skip_tags:
                self._skip_depth = max(0, self._skip_depth - 1)

        def handle_data(self, data):
            if not self._skip_depth and data.strip():
                self._parts.append(data.strip())

        def get_text(self) -> str:
            return "\n".join(self._parts)

    html_text = content.decode("utf-8", errors="replace")
    parser = _TextExtractor()
    parser.feed(html_text)
    text = re.sub(r"\n{3,}", "\n\n", parser.get_text())
    return text.strip()


def _assert_image_magic(content: bytes, ext: str) -> None:
    """Verify image magic bytes to prevent content-type spoofing."""
    if ext in (".png",):
        if not content[:8] == _MAGIC_PNG:
            raise ValueError(f"File does not appear to be a valid PNG (magic bytes mismatch).")
    elif ext in (".jpg", ".jpeg"):
        if not content[:3] == _MAGIC_JPEG:
            raise ValueError(f"File does not appear to be a valid JPEG (magic bytes mismatch).")
    elif ext in (".tiff", ".tif"):
        if content[:4] not in (_MAGIC_TIFF_LE, _MAGIC_TIFF_BE):
            raise ValueError(f"File does not appear to be a valid TIFF (magic bytes mismatch).")


def _ocr_image(content: bytes, _ext: str = "") -> str:
    _assert_image_magic(content, _ext)
    try:
        import pytesseract
        from PIL import Image
        Image.MAX_IMAGE_PIXELS = _MAX_IMAGE_PIXELS
        img = Image.open(io.BytesIO(content))
    except ImportError as exc:
        raise ValueError(
            "OCR not available: install Pillow and pytesseract, and ensure Tesseract is installed."
        ) from exc
    except Image.DecompressionBombError as exc:
        raise ValueError(f"Image size exceeds safe limit: {exc}") from exc
    except Exception as exc:
        raise ValueError(f"not a valid image: {exc}") from exc

    try:
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as exc:
        raise ValueError(f"OCR failed: {exc}") from exc
    finally:
        img.close()  # explicitly free PIL image memory


def _from_text(content: bytes) -> str:
    try:
        return content.decode("utf-8").strip()
    except UnicodeDecodeError:
        raise ValueError(
            "File must be a PDF, DOCX, XLSX, PPTX, HTML, image, or UTF-8 encoded text."
        )
